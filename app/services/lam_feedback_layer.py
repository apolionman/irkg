import faiss
import fitz 
import numpy as np
import json
import os
from bs4 import BeautifulSoup
from pptx import Presentation
import pandas as pd
import docx


from langchain.embeddings.openai import OpenAIEmbeddings

# Initialize OpenAI Embeddings
embeddings = OpenAIEmbeddings()

# FAISS Index File
INDEX_FILE = "/home/dgx/dgx_irkg_be/feedback/feedback_index.faiss"
DATA_FILE = "/home/dgx/dgx_irkg_be/feedback/feedback_data.json"

# Create FAISS index (dimension = 1536 for OpenAI embeddings)
dimension = 1536
index = faiss.IndexFlatL2(dimension)

# Load previous stored data
document_data = {}

if os.path.exists(INDEX_FILE):
    faiss.read_index(INDEX_FILE)
if os.path.exists(DATA_FILE):
    with open(DATA_FILE, "r") as f:
        document_data = json.load(f)

def extract_text_from_file(file_path):
    """ Extracts text from various file types """
    file_ext = file_path.split(".")[-1].lower()
    text = ""

    if file_ext == "pdf":
        doc = fitz.open(file_path)
        text = "\n".join([page.get_text("text") for page in doc])

    elif file_ext in ["doc", "docx"]:
        doc = docx.Document(file_path)
        text = "\n".join([para.text for para in doc.paragraphs])

    elif file_ext in ["xls", "xlsx", "csv"]:
        df = pd.read_excel(file_path) if file_ext in ["xls", "xlsx"] else pd.read_csv(file_path)
        text = df.to_string()

    elif file_ext in ["ppt", "pptx"]:
        prs = Presentation(file_path)
        text = "\n".join([slide.notes_slide.notes_text_frame.text for slide in prs.slides if slide.notes_slide])

    elif file_ext in ["txt"]:
        with open(file_path, "r", encoding="utf-8") as f:
            text = f.read()

    elif file_ext in ["html", "htm", "xml"]:
        with open(file_path, "r", encoding="utf-8") as f:
            soup = BeautifulSoup(f, "lxml")
            text = soup.get_text()

    return text.strip()

def chunk_text(text, chunk_size=500):
    """ Splits text into smaller chunks for vectorization """
    words = text.split()
    return [" ".join(words[i:i+chunk_size]) for i in range(0, len(words), chunk_size)]

def store_file_in_faiss(file_path):
    """ Extracts text, converts to embeddings, and stores in FAISS """
    
    text = extract_text_from_file(file_path)
    if not text:
        return "No text found in the document."

    chunks = chunk_text(text)

    for i, chunk in enumerate(chunks):
        vector = embeddings.embed_query(chunk)
        vector_np = np.array([vector], dtype=np.float32)
        index.add(vector_np)

        document_data[len(document_data)] = chunk

    # Save FAISS index and metadata
    faiss.write_index(index, INDEX_FILE)
    with open(DATA_FILE, "w") as f:
        json.dump(document_data, f)

    return f"Stored {len(chunks)} chunks from {file_path}"

def retrieve_relevant_context(query: str, top_k=3):
    """ Searches FAISS for similar document text chunks """
    
    query_vector = np.array([embeddings.embed_query(query)], dtype=np.float32)

    # Search FAISS for similar vectors
    _, matches = index.search(query_vector, top_k)

    # Retrieve corresponding text chunks
    relevant_contexts = [document_data.get(i, "No relevant context") for i in matches[0] if i in document_data]

    return relevant_contexts


